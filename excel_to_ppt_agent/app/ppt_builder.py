from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from .insights import AnalysisResult, ChartSpec


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, item in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0


def _chart_type_to_xl(chart_type: str) -> XL_CHART_TYPE:
    return {
        "line": XL_CHART_TYPE.LINE,
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "clustered_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    }[chart_type]


def _add_chart_slide(prs: Presentation, spec: ChartSpec) -> None:
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = spec.title

    chart_data = CategoryChartData()
    chart_data.categories = spec.category_labels
    for series_name, values in spec.series.items():
        chart_data.add_series(series_name, values)

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4.5)
    chart_type = _chart_type_to_xl(spec.chart_type)
    chart_shape = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart

    if spec.x_axis_title:
        chart_shape.category_axis.has_title = True
        chart_shape.category_axis.axis_title.text_frame.text = spec.x_axis_title
    if spec.y_axis_title:
        chart_shape.value_axis.has_title = True
        chart_shape.value_axis.axis_title.text_frame.text = spec.y_axis_title


def build_ppt_from_analysis(analysis: AnalysisResult, output_path: Optional[str] = None) -> str:
    prs = Presentation()

    # Title slide
    _add_title_slide(
        prs,
        title="Automated Analysis Report",
        subtitle=Path(analysis.source_path).name,
    )

    # Insights slide
    _add_bullets_slide(prs, title="Key Insights", bullets=analysis.key_insights)

    # Charts
    for chart in analysis.charts:
        _add_chart_slide(prs, chart)

    # Optional: summary stats as a table in a slide
    if analysis.summary_stats:
        cols = ["count", "mean", "std", "min", "max", "sum"]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Summary Statistics"
        num_rows = len(analysis.summary_stats)
        num_cols = 1 + len(cols)
        table = slide.shapes.add_table(num_rows + 1, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)).table
        table.cell(0, 0).text = "Metric"
        for j, col_name in enumerate(cols, start=1):
            table.cell(0, j).text = col_name
        for i, (metric, stat) in enumerate(analysis.summary_stats.items(), start=1):
            table.cell(i, 0).text = metric
            for j, col_name in enumerate(cols, start=1):
                val = stat.get(col_name)
                table.cell(i, j).text = f"{val:,.2f}" if isinstance(val, (int, float)) else ""

    outputs_dir = Path("/workspace/excel_to_ppt_agent/outputs")
    outputs_dir.mkdir(parents=True, exist_ok=True)
    if output_path is None:
        output_path = str(outputs_dir / "report.pptx")
    prs.save(output_path)
    return output_path

